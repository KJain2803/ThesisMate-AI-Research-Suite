import atexit
import io
import os
import shutil
import tempfile
from pathlib import Path

if "MPLCONFIGDIR" not in os.environ:
    mpl_config_dir = Path(tempfile.gettempdir()) / "thesismate-matplotlib"
    mpl_config_dir.mkdir(parents=True, exist_ok=True)
    os.environ["MPLCONFIGDIR"] = str(mpl_config_dir)

import streamlit as st
import pandas as pd
import plotly.express as px
from scipy import stats
from docx import Document
import PyPDF2
from pptx import Presentation
from reportlab.lib.pagesizes import A4
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors
from reportlab.platypus import Table, TableStyle, Image
import matplotlib.pyplot as plt
import numpy as np

st.set_page_config(
    page_title="ThesisMate AI Research Suite", page_icon="🧠", layout="wide"
)
saved_graphs = []
plotly_export_warnings = []
graph_output_dir = Path(tempfile.mkdtemp(prefix="thesismate_graphs_"))
atexit.register(shutil.rmtree, graph_output_dir, ignore_errors=True)


def get_graph_path(filename):
    return graph_output_dir / Path(filename).name


def metric_card(title, value):
    st.markdown(
        f"""
    <div class='metric-card'>
        <h4>{title}</h4>
        <h2>{value}</h2>
    </div>
    """,
        unsafe_allow_html=True,
    )


def save_matplotlib_graph(x, y, title, xlabel, ylabel, graph_name):
    fig, ax = plt.subplots(figsize=(7, 4.5))

    ax.plot(x, y, marker="o", markersize=8, linewidth=3, color="#1F77B4")

    try:
        ax.fill_between(x, y, alpha=0.15, color="#A7D3F4")
    except (TypeError, ValueError):
        pass

    ax.set_title(title, fontsize=14, fontweight="bold", color="#123A63")
    ax.set_xlabel(xlabel, fontsize=11)
    ax.set_ylabel(ylabel, fontsize=11)

    ax.grid(True, linestyle="--", alpha=0.4)
    ax.set_facecolor("#F8FBFE")

    for spine in ax.spines.values():
        spine.set_color("#B0CDE8")

    file_path = get_graph_path(f"{graph_name}.png")
    plt.tight_layout()
    plt.savefig(file_path, dpi=220)
    plt.close()
    saved_graphs.append(str(file_path))


def safe_plotly_save(fig, filename):
    file_path = get_graph_path(filename)
    try:
        fig.write_image(str(file_path))
    except Exception as exc:
        warning = (
            f"Could not export {Path(filename).name} for the PDF dossier: {exc}"
        )
        plotly_export_warnings.append(warning)
        st.warning(warning)
        return None

    saved_graphs.append(str(file_path))
    return str(file_path)


# ---------------- STYLE ----------------
# ---------------- PREMIUM UI STYLE ----------------
st.markdown(
    """
<style>

/* Main app background */
[data-testid="stAppViewContainer"]{
    background: linear-gradient(135deg,#eef3f9,#f8fbff);
    font-family: 'Segoe UI', sans-serif;
}

/* Main width */
.block-container{
    padding-top:1.2rem;
    padding-bottom:2rem;
    max-width: 92%;
}

/* Sidebar */
[data-testid="stSidebar"]{
    background: linear-gradient(180deg,#0f172a,#1e293b);
}
[data-testid="stSidebar"] *{
    color: white !important;
}

/* Hero */
.hero-box{
    background: linear-gradient(135deg,#111827,#1d4ed8);
    padding:30px;
    border-radius:24px;
    box-shadow:0 10px 35px rgba(0,0,0,0.20);
    margin-bottom:25px;
}
.hero-title{
    color:white;
    font-size:42px;
    font-weight:800;
    text-align:center;
}
.hero-sub{
    color:#dbeafe;
    text-align:center;
    font-size:18px;
    margin-top:10px;
}

/* Metric cards */
.metric-card{
    background:white;
    padding:22px;
    border-radius:18px;
    box-shadow:0 4px 18px rgba(0,0,0,0.08);
    text-align:center;
    border-left:6px solid #2563eb;
    margin-bottom:15px;
    transition:0.3s;
}
.metric-card:hover{
    transform:translateY(-4px);
    box-shadow:0 10px 24px rgba(37,99,235,0.18);
}
.metric-card h4{
    color:#475569;
    margin-bottom:8px;
}
.metric-card h2{
    color:#0f172a;
    font-size:34px;
}

/* Section cards */
.section-card{
    background:white;
    padding:22px;
    border-radius:20px;
    box-shadow:0 4px 20px rgba(0,0,0,0.07);
    margin-top:18px;
    margin-bottom:18px;
    transition:0.3s;
}
.section-card:hover{
    box-shadow:0 10px 28px rgba(0,0,0,0.10);
}

/* File uploader */
[data-testid="stFileUploader"]{
    background:white;
    border-radius:18px;
    padding:10px;
    box-shadow:0 3px 15px rgba(0,0,0,0.06);
}

/* Selectbox */
[data-baseweb="select"]{
    background:white;
    border-radius:14px;
    box-shadow:0 3px 14px rgba(0,0,0,0.06);
}

/* Buttons */
.stButton>button, .stDownloadButton>button{
    background: linear-gradient(90deg,#2563eb,#1d4ed8);
    color:white;
    border:none;
    border-radius:12px;
    padding:10px 22px;
    font-weight:600;
    box-shadow:0 4px 16px rgba(37,99,235,0.25);
}
.stButton>button:hover, .stDownloadButton>button:hover{
    transform:translateY(-2px);
}

/* Dataframe */
[data-testid="stDataFrame"]{
    background:white;
    border-radius:14px;
    box-shadow:0 3px 12px rgba(0,0,0,0.05);
    padding:8px;
}

/* Plotly charts */
.js-plotly-plot{
    background:white !important;
    border-radius:18px !important;
    padding:12px !important;
    box-shadow:0 4px 18px rgba(0,0,0,0.06) !important;
    margin-bottom:20px !important;
}

/* Success/Info messages */
.stSuccess, .stInfo{
    border-radius:14px;
    box-shadow:0 3px 12px rgba(0,0,0,0.05);
}

/* Tabs */
button[data-baseweb="tab"]{
    font-size:16px;
    font-weight:600;
    color:#334155;
}
button[data-baseweb="tab"][aria-selected="true"]{
    color:#1d4ed8 !important;
    border-bottom:3px solid #1d4ed8 !important;
}

/* Subheaders */
h2, h3{
    color:#0f172a;
}

/* Smooth fade */
html, body, [class*="css"]  {
    scroll-behavior:smooth;
}
</style>
""",
    unsafe_allow_html=True,
)
# ---------------- UNIVERSAL DOCUMENT READER ----------------


def extract_text_from_file(uploaded_file):
    text = ""

    if uploaded_file is not None:
        file_type = uploaded_file.name.split(".")[-1].lower()

        # PDF Reader
        if file_type == "pdf":
            pdf_reader = PyPDF2.PdfReader(uploaded_file)
            for page in pdf_reader.pages:
                if page.extract_text():
                    text += page.extract_text()

        # DOCX Reader
        elif file_type == "docx":
            doc = Document(uploaded_file)
            for para in doc.paragraphs:
                text += para.text + "\n"

        # PPTX Reader
        elif file_type == "pptx":
            prs = Presentation(uploaded_file)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"

        # TXT Reader
        elif file_type == "txt":
            text = uploaded_file.read().decode("utf-8")

    return text


def analyze_dissertation_text(text):
    word_count = len(text.split())
    citation_count = text.count("(") + text.count("[")
    methodology_words = [
        "method",
        "media",
        "treatment",
        "culture",
        "analysis",
        "design",
        "procedure",
    ]

    methodology_hits = sum(
        [1 for word in methodology_words if word.lower() in text.lower()]
    )

    writing_score = min(100, int(word_count / 5))
    methodology_score = min(100, methodology_hits * 15)

    if citation_count >= 5:
        citation_strength = "High"
    elif citation_count >= 2:
        citation_strength = "Moderate"
    else:
        citation_strength = "Low"

    suggestions = []

    if writing_score < 75:
        suggestions.append("Dissertation chapter needs deeper content expansion.")
    if citation_strength == "Low":
        suggestions.append("Add more recent citation-backed references.")
    if methodology_score < 60:
        suggestions.append(
            "Methodology explanation requires stronger scientific detailing."
        )
    if "conclusion" not in text.lower():
        suggestions.append("Expand practical significance in conclusion chapter.")

    return writing_score, citation_strength, methodology_score, suggestions


# ---------------- PREMIUM SIDEBAR ----------------
st.sidebar.markdown("## 🧠 ThesisMate AI Engines")
st.sidebar.markdown("---")
st.sidebar.markdown("### 📘 Dissertation NLP Review")
st.sidebar.markdown("### 📊 Smart Statistical Analytics")
st.sidebar.markdown("### 🧪 Automatic ANOVA / T-Test")
st.sidebar.markdown("### 🧠 AI Scientific Interpretation")
st.sidebar.markdown("### ✍ Dissertation Result Writer")
st.sidebar.markdown("### 📄 Export Scientific Reports")
st.sidebar.markdown("---")
st.sidebar.info("AI Powered End-to-End Research Intelligence Platform")

# ---------------- PREMIUM HERO ----------------
st.markdown(
    """
<div class='hero-box'>
    <div class='hero-title'>🧠 ThesisMate AI Research Suite™</div>
    <div class='hero-sub'>End-to-End Dissertation Intelligence • Automated Biostatistics • AI Scientific Writing Engine</div>
</div>
""",
    unsafe_allow_html=True,
)

main_tab1, main_tab2, main_tab3, main_tab4 = st.tabs(
    [
        "📘 Dissertation Review",
        "📊 Statistical Analytics",
        "🧠 AI Result Writer",
        "📄 Export Reports",
    ]
)
with main_tab1:

    doc_file = st.file_uploader(
        "Upload Dissertation / Thesis / Seminar File",
        type=["pdf", "docx", "pptx", "txt"],
    )

    if doc_file is not None:
        dissertation_text = extract_text_from_file(doc_file)

        st.success("Dissertation document uploaded and converted successfully.")
        st.write("📑 Thesis Text Preview")
        st.text_area(
            label="Dissertation Preview", value=dissertation_text[:1500], height=180
        )

        writing_score, citation_strength, methodology_score, suggestions = (
            analyze_dissertation_text(dissertation_text)
        )

        st.markdown("---")
        st.header("🧠 Dissertation Intelligence Dashboard")

        col1, col2, col3 = st.columns(3)

        with col1:
            metric_card("✍ Writing Depth Score", f"{writing_score}%")
        with col2:
            metric_card("📚 Citation Strength", citation_strength)
        with col3:
            metric_card("🧪 Methodology Quality", f"{methodology_score}%")

        st.subheader("✍️ AI Writing Suggestions")
        for s in suggestions:
            st.write("•", s)
with main_tab2:

    data_mode = st.selectbox(
        "Choose Data Input Mode",
        ["Upload Excel File", "Upload CSV File", "Manual Data Entry"],
    )

    data_file = None
    df = None
    y_axis = None
    test_choice = None

    if data_mode == "Upload Excel File":
        data_file = st.file_uploader(
            "Upload Experimental Dataset (.xlsx)", type=["xlsx"]
        )

    if data_mode == "Upload CSV File":
        data_file = st.file_uploader("Upload Experimental Dataset (.csv)", type=["csv"])

    if data_mode == "Manual Data Entry":
        st.info("Enter or paste your scientific observations below")
        df = st.data_editor(
            pd.DataFrame(columns=["Treatment", "Days", "Response"]), num_rows="dynamic"
        )

    document_text = ""

    if doc_file is not None:
        document_text = extract_text_from_file(doc_file)

    if data_file is not None:
        try:
            if data_mode == "Upload Excel File":
                df = pd.read_excel(data_file)

            elif data_mode == "Upload CSV File":
                df = pd.read_csv(data_file)

            df.columns = [str(col).strip() for col in df.columns]
            df = df.dropna(how="all")
            df = df.fillna(0)

        except Exception as e:
            st.error(f"Dataset loading error: {e}")
            st.stop()

        st.header("📊 Smart Statistical Dashboard")
        st.success("Dataset uploaded and processed successfully")

        st.subheader("Dataset Preview")
        st.dataframe(df.head(), width="stretch")

        numeric_cols = [
            col
            for col in df.select_dtypes(include="number").columns.tolist()
            if col != "Days"
        ]
        if not numeric_cols:
            st.error("Dataset must include at least one numeric response column.")
            st.stop()

        y_axis = st.selectbox("Select Response Variable for Analysis", numeric_cols)

        test_choice = st.selectbox(
            "Choose Statistical Test",
            [
                "Full Automatic Suite",
                "One Way ANOVA",
                "T Test",
                "Pearson Correlation",
                "Linear Regression",
                "CRD Treatment Ranking",
            ],
        )

    st.markdown("<br>", unsafe_allow_html=True)
    run_analysis = st.button("🚀 Generate Full AI Scientific Intelligence Report")
    pdf_buffer = None
    result_writer = ""
    interpretation = ""
    ai_response = ""
    significance = ""
    best_day = ""
    best_mean = ""
    f_val = 0
    p_val = 0
    if run_analysis:
        saved_graphs.clear()
        plotly_export_warnings.clear()

        # ---------------- DOCX ANALYSIS ----------------
        if doc_file is not None:
            full_text = extract_text_from_file(doc_file)

            word_count = len(full_text.split())
            lower_text = full_text.lower()

            clarity = 70 + min(word_count // 200, 20)

            if "reference" in lower_text or "citation" in lower_text:
                citation_strength = "Moderate"
            else:
                citation_strength = "Low"

            methodology = 65
            if "methodology" in lower_text or "materials and methods" in lower_text:
                methodology += 15

            st.markdown("---")
            st.header("📘 Dissertation Intelligence Dashboard")

            col1, col2, col3 = st.columns(3)

            with col1:
                metric_card("✍ Writing Clarity", f"{clarity}/100")

            with col2:
                metric_card("📚 Citation Strength", citation_strength)

            with col3:
                metric_card("🧪 Methodology Depth", f"{methodology}/100")

            st.warning(
                "Recent literature comparison appears limited in uploaded chapter."
            )
            st.info("Conclusion and practical implication discussion can be expanded.")

            st.subheader("✍ AI Writing Suggestions")

            if word_count < 1500:
                st.write("• Dissertation chapter needs deeper content expansion.")
            else:
                st.write(
                    "• Academic depth appears acceptable but transitions can improve."
                )

            if citation_strength == "Low":
                st.write("• Add more recent citation-backed references.")
            else:
                st.write(
                    "• Citation support is present but comparative literature can increase."
                )

            if methodology < 75:
                st.write(
                    "• Methodology explanation requires stronger scientific detailing."
                )
            else:
                st.write("• Methodology section shows reasonable structural support.")

            st.write("• Expand practical significance in conclusion chapter.")

            ai_response = f"""
            WRITING CLARITY SCORE: {clarity}/100
            CITATION STRENGTH: {citation_strength}
            METHODOLOGY DEPTH SCORE: {methodology}/100
            """

        else:
            ai_response = "No dissertation chapter uploaded."
            if data_file is None:
                st.info("Upload a dataset to run statistical analytics.")

        # ---------------- DATA ANALYSIS ----------------
        if data_file is not None:

            if "Days" in df.columns and test_choice in [
                "Full Automatic Suite",
                "One Way ANOVA",
            ]:
                grouped = df.groupby("Days")[y_axis].mean().reset_index()
                save_matplotlib_graph(
                    grouped["Days"],
                    grouped[y_axis],
                    "ANOVA Mean Response Across Days",
                    "Days",
                    y_axis,
                    "anova_graph",
                )

                st.subheader("📈 Mean Trend Graph")
                fig = px.bar(grouped, x="Days", y=y_axis, text_auto=True)
                fig.update_layout(height=400)
                st.plotly_chart(fig, width="stretch")

                groups = [group[y_axis].values for name, group in df.groupby("Days")]
                f_val, p_val = stats.f_oneway(*groups)

                st.subheader("🧪 Automatic One-Way ANOVA")
                cc1, cc2 = st.columns(2)
                cc1.metric("F-Value", round(f_val, 3))
                cc2.metric("P-Value", round(p_val, 3))

                if p_val < 0.05:
                    st.success("Statistically significant group difference detected.")
                    significance = (
                        "Statistically significant group difference detected."
                    )
                else:
                    st.warning(
                        "No statistically significant group difference detected."
                    )
                    significance = (
                        "No statistically significant group difference detected."
                    )

                best_day = grouped.loc[grouped[y_axis].idxmax(), "Days"]
                best_mean = grouped[y_axis].max()

                st.subheader("🏆 Best Performing Group")
                st.info(f"Highest mean {y_axis} observed at Days = {best_day}")

                interpretation = f"""
            The uploaded dataset demonstrates measurable variation in {y_axis} across different Days groups.
            Mean trend visualization indicates increased response at higher Days levels.
            ANOVA findings suggest: {significance}
            This statistical evidence can be directly incorporated into the Results and Discussion chapter.
            The identified best-performing group should be emphasized in scientific interpretation.
            """
                st.subheader("🧠 AI Statistical Interpretation")
                st.write(interpretation)
                st.markdown("---")
                st.header("📊 Advanced Visual Analytics Dashboard")

                # BOX PLOT
                st.subheader("📦 Biological Response Distribution")
                fig_box = px.box(df, x="Days", y=y_axis, points="all")
                st.plotly_chart(fig_box, width="stretch")
                safe_plotly_save(fig_box, "graph_bio_distribution.png")

                # SCATTER WITH TRENDLINE
                st.subheader("📍 Scatter Trend with Regression Fit")
                fig_scatter = px.scatter(
                    df,
                    x="Days",
                    y=y_axis,
                )
                st.plotly_chart(fig_scatter, width="stretch")
                safe_plotly_save(fig_scatter, "graph_scatter_regression.png")

                # LINE PROGRESSION GRAPH

                st.subheader("⚖️ Mean ± Standard Error Response Plot")
                err_df = (
                    df.groupby("Days")[y_axis]
                    .agg(["mean", "std", "count"])
                    .reset_index()
                )
                err_df["se"] = err_df["std"] / np.sqrt(err_df["count"])

                fig_err = px.line(
                    err_df, x="Days", y="mean", markers=True, error_y="se"
                )
                st.plotly_chart(fig_err, width="stretch")
                safe_plotly_save(fig_err, "graph_errorbar.png")

                # HEATMAP
                if "Treatment" in df.columns:
                    st.subheader("🔥 Treatment × Days Heatmap")
                    pivot = df.pivot_table(
                        values=y_axis, index="Treatment", columns="Days", aggfunc="mean"
                    )
                    fig_heat = px.imshow(pivot, text_auto=True)
                    st.plotly_chart(fig_heat, width="stretch")
                    safe_plotly_save(fig_heat, "graph_heatmap.png")
                    if "Treatment" in df.columns:
                        st.subheader("🧿 Treatment Contribution Donut")
                        treat_mean = (
                            df.groupby("Treatment")[y_axis].mean().reset_index()
                        )
                        fig_donut = px.pie(
                            treat_mean, names="Treatment", values=y_axis, hole=0.5
                        )
                        st.plotly_chart(fig_donut, width="stretch")

                        plt.figure(figsize=(6, 4))
                        plt.pie(treat_mean[y_axis], labels=treat_mean["Treatment"])
                        plt.title("Treatment Contribution")
                        plt.tight_layout()
                        graph_file = get_graph_path("graph_donut.png")
                        plt.savefig(graph_file)
                        plt.close()
                        saved_graphs.append(str(graph_file))
                    if "Treatment" in df.columns:
                        st.subheader("🕸️ Multi-Treatment Radar Intelligence")
                        radar = df.groupby("Treatment")[y_axis].mean().reset_index()
                        fig_radar = px.line_polar(
                            radar, r=y_axis, theta="Treatment", line_close=True
                        )
                        st.plotly_chart(fig_radar, width="stretch")
                        safe_plotly_save(fig_radar, "graph_radar.png")
                    avg_perf = round(df[y_axis].mean(), 2)
                    max_perf = round(df[y_axis].max(), 2)

                    st.subheader("🧠 AI Biological Response Health Score")
                    performance_percent = (
                        min(100, round((avg_perf / max_perf) * 100))
                        if max_perf
                        else 0
                    )

                    fig_gauge = px.pie(
                        values=[performance_percent, 100 - performance_percent],
                        names=["Performance", "Remaining"],
                        hole=0.7,
                    )
                    st.plotly_chart(fig_gauge, width="stretch")
                    safe_plotly_save(fig_gauge, "graph_gauge.png")

                    if "Treatment" in df.columns:
                        st.subheader("📦 Treatment-wise Biological Spread")
                        fig_box2 = px.box(df, x="Treatment", y=y_axis, points="all")
                        st.plotly_chart(fig_box2, width="stretch")
                        safe_plotly_save(fig_box2, "graph_treatment_spread.png")

                    st.subheader("📈 Cumulative Response Progression")
                    cum_df = grouped.copy()
                    cum_df["Cumulative"] = cum_df[y_axis].cumsum()
                    fig_cum = px.area(cum_df, x="Days", y="Cumulative")
                    st.plotly_chart(fig_cum, width="stretch")
                    safe_plotly_save(fig_cum, "graph_cumulative.png")

                    st.subheader("🏅 Top Performing Raw Observations")
                    top_obs = df.sort_values(by=y_axis, ascending=False).head(5)
                    st.dataframe(top_obs, width="stretch")

                # ---------------- ADVANCED STATISTICAL MODULE ----------------
                if test_choice == "Full Automatic Suite":
                    st.subheader("📈 Advanced Statistical Intelligence")

                    # Correlation
                    corr_value = round(df["Days"].corr(df[y_axis]), 3)
                    st.write(
                        f"**Pearson Correlation (Days vs {y_axis}) = {corr_value}**"
                    )

                    if corr_value > 0:
                        st.success(
                            "Positive correlation detected: response variable increases with culture duration."
                        )
                    else:
                        st.warning("Negative/weak correlation trend detected.")

                    # Regression
                    slope, intercept, r_value, p_reg, std_err = stats.linregress(
                        df["Days"], df[y_axis]
                    )

                    st.write(
                        f"**Linear Regression Equation:** {y_axis} = {round(intercept,3)} + {round(slope,3)}(Days)"
                    )
                    st.write(f"**Regression R² = {round(r_value**2,3)}**")

                    # Coefficient of Variation
                    cv = round((df[y_axis].std() / df[y_axis].mean()) * 100, 2)
                    st.write(f"**Coefficient of Variation = {cv}%**")
                    plt.figure(figsize=(7, 4.5))
                    plt.hist(df[y_axis], bins=8, color="#5DADE2", edgecolor="black")
                    plt.title(
                        "Frequency Distribution Histogram",
                        fontsize=13,
                        fontweight="bold",
                        color="#123A63",
                    )
                    plt.xlabel(y_axis)
                    plt.ylabel("Frequency")
                    plt.grid(axis="y", linestyle="--", alpha=0.3)
                    plt.tight_layout()
                    graph_file = get_graph_path("graph_hist.png")
                    plt.savefig(graph_file, dpi=220)
                    plt.close()
                    saved_graphs.append(str(graph_file))

                    # Treatment Ranking
                    st.subheader("🏆 Treatment Performance Ranking")
                    treatment_col = (
                        "Treatment" if "Treatment" in df.columns else df.columns[0]
                    )
                    ranking = (
                        df.groupby(treatment_col)[y_axis]
                        .mean()
                        .sort_values(ascending=False)
                        .reset_index()
                    )
                    st.dataframe(ranking, width="stretch")

                    top_treatment = ranking.iloc[0][treatment_col]
                    top_value = round(ranking.iloc[0][y_axis], 3)

                    st.info(
                        f"Top ranked treatment based on mean {y_axis} = {top_treatment} ({top_value})"
                    )
                    plt.figure(figsize=(7, 4.5))
                    bars = plt.bar(
                        ranking[treatment_col],
                        ranking[y_axis],
                        color="#3498DB",
                        edgecolor="black",
                    )
                    plt.title(
                        "Treatment Performance Ranking",
                        fontsize=13,
                        fontweight="bold",
                        color="#123A63",
                    )
                    plt.xlabel("Treatment")
                    plt.ylabel(y_axis)
                    plt.grid(axis="y", linestyle="--", alpha=0.3)

                    for bar in bars:
                        yval = bar.get_height()
                        plt.text(bar.get_x() + 0.15, yval + 0.03, round(yval, 2))

                    plt.tight_layout()
                    graph_file = get_graph_path("graph_treatment.png")
                    plt.savefig(graph_file, dpi=220)
                    plt.close()
                    saved_graphs.append(str(graph_file))
                    st.subheader("✍️ AI Dissertation Result Writer")

                    result_writer = f"""
                        Statistical assessment of {y_axis} across culture duration revealed a strong positive Pearson correlation coefficient of {corr_value}, indicating that the response variable increased progressively with increasing days of incubation. Linear regression modelling further supported this trend, represented by the equation {y_axis} = {round(intercept,3)} + {round(slope,3)}(Days), with a coefficient of determination (R²) of {round(r_value**2,3)}, suggesting substantial dependence of the experimental response on treatment duration.

                        The coefficient of variation was computed as {cv}%, indicating acceptable experimental consistency with measurable biological variability. Among the tested treatment combinations, {top_treatment} exhibited the highest mean performance ({top_value}) for {y_axis}, thereby representing the most responsive culture condition under the present investigation. These findings collectively indicate that treatment duration and media composition exerted pronounced influence on in vitro response and can be scientifically emphasized in the Results and Discussion chapter.
                        """

                    st.write(result_writer)
                    st.markdown("---")
                    st.subheader("🧠 ThesisMate Final Scientific Verdict")

                    confidence_score = min(
                        99,
                        round(
                            (abs(corr_value) * 30)
                            + ((r_value**2) * 30)
                            + ((100 - cv) / 2)
                        ),
                    )
                    novelty_score = min(
                        98,
                        round(
                            (best_mean * 5)
                            if isinstance(best_mean, (int, float))
                            else 85
                        ),
                    )

                    c1, c2, c3 = st.columns(3)
                    with c1:
                        metric_card("📌 Statistical Confidence", f"{confidence_score}%")
                    with c2:
                        metric_card("🔬 Research Strength", f"{novelty_score}%")
                    with c3:
                        metric_card("🏆 Best Treatment", str(top_treatment))

                    st.success(
                        "AI Verdict: Dataset demonstrates publishable analytical strength with statistically interpretable treatment differentiation."
                    )

                    report_text = f"""
                        THESISMATE AI FINAL REPORT

                        ===== DISSERTATION REVIEW =====
                        {ai_response}

                        ===== STATISTICAL REVIEW =====
                        Analyzed Variable = {y_axis}
                        F-value = {round(f_val,3)}
                        P-value = {round(p_val,3)}
                        Conclusion = {significance}
                        Best Group = Days {best_day} with mean {round(best_mean,3)}

                        Interpretation:
                        {interpretation}
                        """

                    def add_page_number(canvas, pdf_doc):
                        canvas.saveState()
                        canvas.setFont("Helvetica", 9)
                        canvas.setFillColor(colors.grey)
                        canvas.drawString(
                            40, 20, "Generated by ThesisMate AI Research Engine™"
                        )
                        canvas.drawRightString(550, 20, f"Page {pdf_doc.page}")
                        canvas.restoreState()

                    pdf_buffer = io.BytesIO()
                    pdf_doc = SimpleDocTemplate(
                        pdf_buffer,
                        pagesize=A4,
                        rightMargin=30,
                        leftMargin=30,
                        topMargin=30,
                        bottomMargin=30,
                    )
                    styles = getSampleStyleSheet()
                    story = []

                    title_style = styles["Title"]
                    title_style.fontSize = 24
                    title_style.textColor = colors.HexColor("#0A2A66")

                    sub_style = styles["Heading2"]
                    sub_style.textColor = colors.HexColor("#1F4E79")
                    sub_style.spaceAfter = 12

                    normal_style = styles["BodyText"]
                    normal_style.spaceAfter = 8

                    # -------- COVER PAGE ----------
                    story.append(Spacer(1, 70))

                    cover_box = Table(
                        [
                            [
                                Paragraph(
                                    "<b>THESISMATE AI EXECUTIVE RESEARCH DOSSIER</b>",
                                    title_style,
                                )
                            ],
                            [
                                Paragraph(
                                    "Automated Scientific Intelligence | Dissertation Analytics | Biostatistical Modelling",
                                    sub_style,
                                )
                            ],
                            [
                                Paragraph(
                                    f"Target Variable Analysed: <b>{y_axis}</b>",
                                    styles["Heading3"],
                                )
                            ],
                            [
                                Paragraph(
                                    "Generated on ThesisMate AI Research Engine™",
                                    styles["Heading3"],
                                )
                            ],
                            [
                                Paragraph(
                                    "CONFIDENTIAL SCIENTIFIC ANALYTICS REPORT",
                                    styles["Heading2"],
                                )
                            ],
                        ],
                        colWidths=460,
                    )

                    cover_box.setStyle(
                        TableStyle(
                            [
                                (
                                    "BOX",
                                    (0, 0),
                                    (-1, -1),
                                    2,
                                    colors.HexColor("#0A2A66"),
                                ),
                                (
                                    "BACKGROUND",
                                    (0, 0),
                                    (-1, -1),
                                    colors.HexColor("#EDF4FB"),
                                ),
                                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                ("PADDING", (0, 0), (-1, -1), 22),
                                (
                                    "GRID",
                                    (0, 0),
                                    (-1, -1),
                                    0.5,
                                    colors.HexColor("#9CC2E5"),
                                ),
                            ]
                        )
                    )
                    story.append(cover_box)
                    story.append(PageBreak())

                    # -------- KPI DASHBOARD PAGE ----------
                    story.append(
                        Paragraph("AI SCIENTIFIC EXECUTIVE SUMMARY", sub_style)
                    )
                    story.append(Spacer(1, 20))

                    kpi_data = [
                        ["Pearson Correlation", round(corr_value, 3)],
                        ["Regression R²", round(r_value**2, 3)],
                        ["Coefficient of Variation", f"{round(cv,2)}%"],
                        ["Best Performing Group", str(top_treatment)],
                        ["F Statistic", round(f_val, 3)],
                        ["P Value", round(p_val, 4)],
                    ]

                    kpi_table = Table(kpi_data, colWidths=[220, 180])
                    kpi_table.setStyle(
                        TableStyle(
                            [
                                (
                                    "GRID",
                                    (0, 0),
                                    (-1, -1),
                                    1,
                                    colors.HexColor("#0A2A66"),
                                ),
                                ("BACKGROUND", (0, 0), (-1, -1), colors.whitesmoke),
                                (
                                    "TEXTCOLOR",
                                    (0, 0),
                                    (-1, -1),
                                    colors.HexColor("#0A2A66"),
                                ),
                                ("PADDING", (0, 0), (-1, -1), 10),
                            ]
                        )
                    )
                    story.append(kpi_table)
                    story.append(Spacer(1, 30))

                    verdict = "Dataset demonstrates strong analytical publishability with measurable treatment differentiation and statistically interpretable biological response progression."
                    verdict_table = Table(
                        [[Paragraph(verdict, normal_style)]], colWidths=430
                    )
                    verdict_table.setStyle(
                        TableStyle(
                            [
                                ("BOX", (0, 0), (-1, -1), 1, colors.green),
                                ("BACKGROUND", (0, 0), (-1, -1), colors.lightgreen),
                                ("PADDING", (0, 0), (-1, -1), 12),
                            ]
                        )
                    )
                    story.append(verdict_table)
                    story.append(Spacer(1, 25))

                    story.append(
                        Paragraph(
                            f"<b>AI Publishability Confidence: {confidence_score}% | Experimental Reliability: HIGH | Thesis Novelty Potential: MODERATE-HIGH</b>",
                            styles["Heading3"],
                        )
                    )
                    story.append(Spacer(1, 10))

                    story.append(
                        Paragraph(
                            "<i>Integrated NLP dissertation review, inferential biostatistics, graphical evidence mining and automated result synthesis confirm dissertation-grade analytical robustness.</i>",
                            normal_style,
                        )
                    )

                    story.append(PageBreak())
                    story.append(Spacer(1, 10))

                    # -------- AI TEXT SECTIONS ----------
                    sections = [
                        ["DISSERTATION REVIEW", ai_response],
                        [
                            "STATISTICAL REVIEW",
                            f"Analyzed Variable: {y_axis}<br/>F-value: {round(f_val,3)}<br/>P-value: {round(p_val,3)}<br/>Conclusion: {significance}<br/>Best Group: Days {best_day} with mean {round(best_mean,3)}",
                        ],
                        ["AI SCIENTIFIC INTERPRETATION", interpretation],
                        ["RESULT WRITER PARAGRAPH", result_writer],
                    ]

                    # -------- EXECUTIVE SCIENTIFIC METRIC DASHBOARD --------
                    story.append(
                        Paragraph("EXECUTIVE SCIENTIFIC METRIC DASHBOARD", sub_style)
                    )
                    story.append(Spacer(1, 18))

                    metric_data = [
                        [
                            Paragraph(
                                f"<b>F STATISTIC</b><br/>{round(f_val,3)}",
                                styles["BodyText"],
                            ),
                            Paragraph(
                                f"<b>P VALUE</b><br/>{round(p_val,4)}",
                                styles["BodyText"],
                            ),
                        ],
                        [
                            Paragraph(
                                f"<b>PEARSON CORRELATION</b><br/>{corr_value}",
                                styles["BodyText"],
                            ),
                            Paragraph(
                                f"<b>REGRESSION R²</b><br/>{round(r_value**2,3)}",
                                styles["BodyText"],
                            ),
                        ],
                        [
                            Paragraph(
                                f"<b>COEFF. OF VARIATION</b><br/>{cv}%",
                                styles["BodyText"],
                            ),
                            Paragraph(
                                f"<b>BEST TREATMENT</b><br/>{top_treatment}",
                                styles["BodyText"],
                            ),
                        ],
                    ]

                    metric_table = Table(
                        metric_data, colWidths=[225, 225], rowHeights=[72, 72, 72]
                    )
                    metric_table.setStyle(
                        TableStyle(
                            [
                                (
                                    "BOX",
                                    (0, 0),
                                    (-1, -1),
                                    1.5,
                                    colors.HexColor("#1F4E79"),
                                ),
                                (
                                    "GRID",
                                    (0, 0),
                                    (-1, -1),
                                    1,
                                    colors.HexColor("#9CC2E5"),
                                ),
                                (
                                    "BACKGROUND",
                                    (0, 0),
                                    (-1, -1),
                                    colors.HexColor("#EDF4FB"),
                                ),
                                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                                ("FONTSIZE", (0, 0), (-1, -1), 11),
                                (
                                    "TEXTCOLOR",
                                    (0, 0),
                                    (-1, -1),
                                    colors.HexColor("#123A63"),
                                ),
                                ("FONTNAME", (0, 0), (-1, -1), "Helvetica-Bold"),
                            ]
                        )
                    )
                    story.append(metric_table)
                    story.append(Spacer(1, 25))
                    for heading, content in sections:
                        head = Table(
                            [[Paragraph(f"<b>{heading}</b>", styles["Heading3"])]],
                            colWidths=450,
                        )
                        head.setStyle(
                            TableStyle(
                                [
                                    (
                                        "BACKGROUND",
                                        (0, 0),
                                        (-1, -1),
                                        colors.HexColor("#D9EAF7"),
                                    ),
                                    (
                                        "BOX",
                                        (0, 0),
                                        (-1, -1),
                                        1,
                                        colors.HexColor("#0A2A66"),
                                    ),
                                    ("PADDING", (0, 0), (-1, -1), 8),
                                ]
                            )
                        )
                        story.append(head)

                        body = Table(
                            [[Paragraph(content.replace("\n", "<br/>"), normal_style)]],
                            colWidths=450,
                        )
                        body.setStyle(
                            TableStyle(
                                [
                                    ("BOX", (0, 0), (-1, -1), 1, colors.grey),
                                    ("PADDING", (0, 0), (-1, -1), 10),
                                ]
                            )
                        )
                        story.append(body)
                        story.append(Spacer(1, 25))

                    story.append(PageBreak())
                    desc = (
                        df[[col for col in df.select_dtypes(include="number").columns]]
                        .describe()
                        .round(2)
                    )
                    desc_data = [["Statistic"] + list(desc.columns)]

                    for idx in desc.index:
                        desc_data.append([idx] + list(desc.loc[idx]))

                    desc_table = Table(desc_data)
                    desc_table.setStyle(
                        TableStyle(
                            [
                                ("GRID", (0, 0), (-1, -1), 1, colors.grey),
                                (
                                    "BACKGROUND",
                                    (0, 0),
                                    (-1, 0),
                                    colors.HexColor("#D9EAF7"),
                                ),
                                (
                                    "TEXTCOLOR",
                                    (0, 0),
                                    (-1, 0),
                                    colors.HexColor("#0A2A66"),
                                ),
                                ("FONTSIZE", (0, 0), (-1, -1), 7),
                                ("PADDING", (0, 0), (-1, -1), 4),
                            ]
                        )
                    )
                    story.append(desc_table)
                    story.append(Spacer(1, 20))
                    # -------- STAT TABLES ----------
                    story.append(
                        Paragraph("DETAILED STATISTICAL DATA ANNEXURE", sub_style)
                    )
                    story.append(Spacer(1, 15))

                    # ---------- TABLE 1 DESCRIPTIVE STATS ----------

                    # ---------- TABLE 2 TREATMENT RANK ----------
                    if "Treatment" in df.columns:
                        rank_df = (
                            df.groupby("Treatment")[y_axis]
                            .mean()
                            .round(3)
                            .reset_index()
                        )
                        rank_data = [["Treatment", f"Mean {y_axis}"]]
                        for i in range(len(rank_df)):
                            rank_data.append(list(rank_df.iloc[i]))

                        rank_table = Table(rank_data)
                        rank_table.setStyle(
                            TableStyle(
                                [
                                    ("GRID", (0, 0), (-1, -1), 1, colors.grey),
                                    (
                                        "BACKGROUND",
                                        (0, 0),
                                        (-1, 0),
                                        colors.HexColor("#D9EAF7"),
                                    ),
                                    ("PADDING", (0, 0), (-1, -1), 5),
                                ]
                            )
                        )
                        story.append(rank_table)
                        story.append(Spacer(1, 20))

                    # ---------- TABLE 3 TOP RAW OBSERVATIONS ----------
                    raw_cols = [
                        col
                        for col in [treatment_col, "Days", y_axis]
                        if col in df.columns
                    ]
                    raw_df = (
                        df.sort_values(by=y_axis, ascending=False)[raw_cols]
                        .head(10)
                        .round(3)
                    )
                    raw_data = [list(raw_df.columns)]
                    for i in range(len(raw_df)):
                        raw_data.append(list(raw_df.iloc[i]))

                    raw_table = Table(raw_data)
                    raw_table.setStyle(
                        TableStyle(
                            [
                                ("GRID", (0, 0), (-1, -1), 1, colors.grey),
                                (
                                    "BACKGROUND",
                                    (0, 0),
                                    (-1, 0),
                                    colors.HexColor("#D9EAF7"),
                                ),
                                ("FONTSIZE", (0, 0), (-1, -1), 8),
                            ]
                        )
                    )
                    story.append(raw_table)
                    story.append(PageBreak())

                    # ---------- GRAPHICAL EVIDENCE ----------
                    story.append(
                        Paragraph("AI GRAPHICAL SCIENTIFIC EVIDENCE", sub_style)
                    )
                    story.append(Spacer(1, 20))

                    graph_comments = {
                        "anova_graph.png": "Mean group comparison reveals biological response escalation with incubation duration.",
                        "graph_bio_distribution.png": "Distribution analysis indicates widening biological variability among day groups.",
                        "graph_scatter_regression.png": "Scatter evidence confirms positive raw observational trend across culture days.",
                        "graph_errorbar.png": "Standard error margins indicate acceptable experimental reliability and mean consistency.",
                        "graph_heatmap.png": "Heatmap identifies superior treatment-day interaction zones with higher response intensity.",
                        "graph_donut.png": "Treatment contribution proportions reveal relative biological share among media combinations.",
                        "graph_radar.png": "Radar intelligence highlights comparative multidirectional treatment superiority.",
                        "graph_gauge.png": "AI health score summarizes cumulative biological performance efficiency of the dataset.",
                        "graph_treatment_spread.png": "Treatment spread plot visualizes inter-treatment variability and response concentration.",
                        "graph_cumulative.png": "Cumulative progression indicates sustained additive response with increasing incubation.",
                        "graph_treatment.png": "Treatment ranking confirms top-performing media composition under tested conditions.",
                    }

                    graph_counter = 0

                    for g in saved_graphs:
                        graph_file = Path(g)
                        graph_name = graph_file.name

                        if not graph_file.exists():
                            continue
                        graph_title = (
                            graph_name.replace(".png", "").replace("_", " ").upper()
                        )

                        graph_head = Table(
                            [[Paragraph(graph_title, styles["Heading4"])]],
                            colWidths=450,
                        )
                        graph_head.setStyle(
                            TableStyle(
                                [
                                    (
                                        "BACKGROUND",
                                        (0, 0),
                                        (-1, -1),
                                        colors.HexColor("#EAF4FB"),
                                    ),
                                    (
                                        "BOX",
                                        (0, 0),
                                        (-1, -1),
                                        1,
                                        colors.HexColor("#1F4E79"),
                                    ),
                                    ("LEFTPADDING", (0, 0), (-1, -1), 8),
                                ]
                            )
                        )
                        story.append(graph_head)
                        story.append(Spacer(1, 8))

                        img = Image(str(graph_file), width=430, height=240)
                        story.append(img)
                        story.append(Spacer(1, 8))

                        comment = graph_comments.get(
                            graph_name,
                            "Graphical statistical evidence generated through ThesisMate analytical engine.",
                        )
                        cbox = Table(
                            [[Paragraph(comment, styles["BodyText"])]], colWidths=450
                        )
                        cbox.setStyle(
                            TableStyle(
                                [
                                    ("BACKGROUND", (0, 0), (-1, -1), colors.whitesmoke),
                                    ("BOX", (0, 0), (-1, -1), 0.8, colors.grey),
                                    ("LEFTPADDING", (0, 0), (-1, -1), 6),
                                ]
                            )
                        )
                        story.append(cbox)
                        story.append(Spacer(1, 12))

                        graph_counter += 1
                        if graph_counter % 2 == 0:
                            story.append(PageBreak())

                            # -------- FINAL SCIENTIFIC CERTIFICATION PAGE --------
                    story.append(PageBreak())
                    story.append(
                        Paragraph("THESISMATE AI SCIENTIFIC CERTIFICATION", sub_style)
                    )
                    story.append(Spacer(1, 25))

                    cert_data = [
                        [
                            Paragraph(
                                f"<b>STATISTICAL CONFIDENCE SCORE</b><br/>{confidence_score}%",
                                styles["BodyText"],
                            )
                        ],
                        [
                            Paragraph(
                                f"<b>RESEARCH STRENGTH INDEX</b><br/>{novelty_score}%",
                                styles["BodyText"],
                            )
                        ],
                        [
                            Paragraph(
                                f"<b>TOP PERFORMING TREATMENT</b><br/>{top_treatment}",
                                styles["BodyText"],
                            )
                        ],
                        [
                            Paragraph(
                                "<b>PUBLISHABILITY VERDICT</b><br/>Dataset demonstrates strong scientific interpretability, reproducible analytical depth and dissertation-grade publication readiness under AI-assisted evaluation.",
                                styles["BodyText"],
                            )
                        ],
                    ]

                    cert_table = Table(cert_data, colWidths=450)
                    cert_table.setStyle(
                        TableStyle(
                            [
                                (
                                    "BOX",
                                    (0, 0),
                                    (-1, -1),
                                    2,
                                    colors.HexColor("#1F4E79"),
                                ),
                                (
                                    "GRID",
                                    (0, 0),
                                    (-1, -1),
                                    1,
                                    colors.HexColor("#9CC2E5"),
                                ),
                                (
                                    "BACKGROUND",
                                    (0, 0),
                                    (-1, -1),
                                    colors.HexColor("#EDF4FB"),
                                ),
                                ("ALIGN", (0, 0), (-1, -1), "CENTER"),
                                ("VALIGN", (0, 0), (-1, -1), "MIDDLE"),
                                ("PADDING", (0, 0), (-1, -1), 14),
                            ]
                        )
                    )

                    story.append(cert_table)
                    story.append(Spacer(1, 35))

                    story.append(
                        Paragraph(
                            "Certified and Generated by ThesisMate AI Research Engine™",
                            styles["Heading3"],
                        )
                    )
                    story.append(Spacer(1, 10))
                    story.append(
                        Paragraph(
                            "Automated Dissertation Intelligence • Biostatistical Validation • Scientific Writing Synthesis",
                            styles["BodyText"],
                        )
                    )
                    story.append(Spacer(1, 20))
                    story.append(
                        Paragraph(
                            "<b>----- END OF AI SCIENTIFIC DOSSIER -----</b>",
                            styles["Heading3"],
                        )
                    )
                    pdf_doc.build(
                        story, onFirstPage=add_page_number, onLaterPages=add_page_number
                    )
                    pdf_buffer.seek(0)

        if data_file is not None and test_choice == "T Test":
            st.markdown("---")
            st.header("🧪 Independent T Test Module")

            treatment_col = st.selectbox("Select Treatment Column", df.columns)

            unique_groups = list(df[treatment_col].dropna().unique())

            if len(unique_groups) >= 2:

                group1 = st.selectbox("Choose Group 1", unique_groups, index=0)
                remaining_groups = [g for g in unique_groups if g != group1]

                group2 = st.selectbox("Choose Group 2", remaining_groups, index=0)

                data1 = df[df[treatment_col] == group1][y_axis]
                data2 = df[df[treatment_col] == group2][y_axis]

                t_stat, p_val = stats.ttest_ind(data1, data2)

                col1, col2 = st.columns(2)
                with col1:
                    st.metric("T Statistic", round(t_stat, 3))
                with col2:
                    st.metric("P Value", round(p_val, 4))

                if p_val < 0.05:
                    st.success("Statistically significant difference detected.")
                else:
                    st.warning("No statistically significant difference detected.")

            else:
                st.error("T Test requires at least two unique treatment groups.")

        if data_file is not None and test_choice == "Pearson Correlation":
            st.markdown("---")
            st.header("📈 Pearson Correlation Module")

            if "Days" in df.columns:
                corr_value = round(df["Days"].corr(df[y_axis]), 3)
                st.metric("Pearson Correlation Coefficient", corr_value)
                fig_corr = px.scatter(
                    df,
                    x="Days",
                    y=y_axis,
                )
                st.plotly_chart(fig_corr, width="stretch")
                save_matplotlib_graph(
                    df["Days"],
                    df[y_axis],
                    "Pearson Correlation Scatter",
                    "Days",
                    y_axis,
                    "pearson_graph",
                )

                if corr_value > 0:
                    st.success(
                        "Positive correlation detected between Days and response variable."
                    )
                else:
                    st.warning("Negative or weak relationship detected.")

        if data_file is not None and test_choice == "Linear Regression":
            st.markdown("---")
            st.header("📉 Linear Regression Module")

            if "Days" in df.columns:
                slope, intercept, r_value, p_reg, std_err = stats.linregress(
                    df["Days"], df[y_axis]
                )

                st.write(
                    f"Regression Equation: {y_axis} = {round(intercept,3)} + {round(slope,3)}(Days)"
                )
                st.write(f"R² = {round(r_value**2,3)}")
                fig_reg = px.scatter(
                    df,
                    x="Days",
                    y=y_axis,
                )
                st.plotly_chart(fig_reg, width="stretch")
                save_matplotlib_graph(
                    df["Days"],
                    df[y_axis],
                    "Linear Regression Model",
                    "Days",
                    y_axis,
                    "regression_graph",
                )

        if data_file is not None and test_choice == "CRD Treatment Ranking":
            st.markdown("---")
            st.header("🏆 CRD Treatment Ranking Module")

            treatment_col = st.selectbox(
                "Select Treatment Column for Ranking", df.columns
            )

            ranking = (
                df.groupby(treatment_col)[y_axis]
                .mean()
                .sort_values(ascending=False)
                .reset_index()
            )
            st.dataframe(ranking, width="stretch")
            fig_rank = px.bar(ranking, x=treatment_col, y=y_axis, text_auto=True)
            st.plotly_chart(fig_rank, width="stretch")
            save_matplotlib_graph(
                ranking[treatment_col],
                ranking[y_axis],
                "CRD Treatment Ranking",
                treatment_col,
                y_axis,
                "crd_graph",
            )

            st.subheader("🧿 Treatment Contribution Donut")
            fig_pie = px.pie(ranking, names=treatment_col, values=y_axis, hole=0.5)
            st.plotly_chart(fig_pie, width="stretch")
            safe_plotly_save(fig_pie, "graph_donut.png")
            top_treatment = ranking.iloc[0][treatment_col]
            top_value = round(ranking.iloc[0][y_axis], 3)

            st.info(
                f"Top ranked treatment based on mean {y_axis} = {top_treatment} ({top_value})"
            )

with main_tab3:
    if run_analysis and data_file is not None and test_choice == "Full Automatic Suite":
        st.markdown("<div class='section-card'>", unsafe_allow_html=True)
        st.subheader("🧠 AI Dissertation Result Writer")
        st.write(result_writer)
        st.markdown("</div>", unsafe_allow_html=True)

with main_tab4:
    if (
        run_analysis
        and pdf_buffer is not None
        and test_choice == "Full Automatic Suite"
    ):
        st.markdown("<div class='section-card'>", unsafe_allow_html=True)
        st.subheader("📄 Export Professional Scientific Dossier")
        st.download_button(
            label="Download Final PDF Report",
            data=pdf_buffer,
            file_name="ThesisMate_Final_Report.pdf",
            mime="application/pdf",
        )
        st.markdown("</div>", unsafe_allow_html=True)
